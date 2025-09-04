import os
import textwrap
import nbformat
import base64
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
from io import BytesIO

NOTEBOOK_PATH = os.path.join(os.path.dirname(__file__), "P09_jupyternotebook.ipynb")
OUTPUT_PPTX = os.path.join(os.path.dirname(__file__), "P09_presentation_with_outputs.pptx")
MAX_CHARS_PER_LINE = 110

def add_markdown_slide(prs, md_text):
    lines = [l.rstrip() for l in md_text.splitlines() if l.strip() != ""]
    title_text = None
    body_text = []
    for i, l in enumerate(lines):
        if l.strip().startswith("#"):
            title_text = l.lstrip("#").strip()
            body_text = lines[i+1:]
            break
    if title_text is None:
        title_text = lines[0] if lines else "Notes"
        body_text = lines[1:]
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for p in body_text:
        if p.startswith("- ") or p.startswith("* "):
            p = p[2:].strip()
            p = textwrap.fill(p, width=MAX_CHARS_PER_LINE)
            p_pr = tf.add_paragraph()
            p_pr.level = 1
            p_pr.text = p
        else:
            p = textwrap.fill(p, width=MAX_CHARS_PER_LINE)
            p_pr = tf.add_paragraph()
            p_pr.level = 0
            p_pr.text = p

def add_code_slide(prs, code_text, title="Code"):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5); top = Inches(0.2); width = Inches(9); height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    t = title_box.text_frame
    t.text = title
    t.paragraphs[0].font.size = Pt(22)
    left = Inches(0.5); top = Inches(1.0); width = prs.slide_width - Inches(1.0); height = prs.slide_height - Inches(1.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_top = Pt(6); tf.margin_right = Pt(6); tf.margin_bottom = Pt(6)
    for i, line in enumerate(code_text.splitlines()):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        p.font.name = 'Courier New'
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.LEFT

def add_output_text_slide(prs, text, title="Output"):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5); top = Inches(0.2); width = Inches(9); height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_box.text_frame.text = title
    left = Inches(0.5); top = Inches(1.0); width = prs.slide_width - Inches(1.0); height = prs.slide_height - Inches(1.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for line in str(text).splitlines():
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = 'Courier New'

def add_output_image_slide(prs, image_bytes, title="Output Image"):
    # save bytes to temp file and insert
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name
    try:
        # open to get size and scale to slide width while preserving aspect ratio
        img = Image.open(tmp_path)
        img_w, img_h = img.size
        slide_w = prs.slide_width
        max_width = slide_w - Inches(1.0)
        # compute width in px equivalent: pptx uses English Metric Units; we use width directly when adding
        left = Inches(0.5)
        top = Inches(1.0)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(left, Inches(0.2), max_width, Inches(0.6))
        title_box.text_frame.text = title
        # add picture scaled to max_width
        slide.shapes.add_picture(tmp_path, left, top, width=max_width)
    finally:
        try:
            os.remove(tmp_path)
        except Exception:
            pass

def notebook_to_pptx(nb_path, out_path):
    nb = nbformat.read(nb_path, as_version=4)
    prs = Presentation()
    title = "Presentation"
    if nb.metadata and nb.metadata.get("title"):
        title = nb.metadata["title"]
    else:
        for cell in nb.cells:
            if cell.cell_type == "markdown":
                for ln in cell.source.splitlines():
                    if ln.strip().startswith("#"):
                        title = ln.lstrip("#").strip()
                        break
                if title != "Presentation":
                    break
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = "Generated from notebook (includes outputs & images)"

    for idx, cell in enumerate(nb.cells):
        if cell.cell_type == "markdown":
            if cell.source.strip():
                add_markdown_slide(prs, cell.source)
        elif cell.cell_type == "code":
            if cell.source.strip():
                first_line = cell.source.strip().splitlines()[0]
                code_title = first_line if len(first_line) < 60 else first_line[:60]
                add_code_slide(prs, cell.source, title=code_title)
            # handle outputs
            outputs = cell.get('outputs', []) or []
            for oidx, out in enumerate(outputs):
                # stream outputs (stdout/stderr)
                if out.get('output_type') == 'stream':
                    text = out.get('text', '')
                    if text:
                        add_output_text_slide(prs, text, title=f"Output (stream) [{idx}.{oidx}]")
                # execution/display results
                elif out.get('output_type') in ('execute_result', 'display_data'):
                    data = out.get('data', {})
                    # image/png or image/jpeg
                    if 'image/png' in data:
                        b = base64.b64decode(data['image/png'])
                        add_output_image_slide(prs, b, title=f"Image [{idx}.{oidx}]")
                    elif 'image/jpeg' in data:
                        b = base64.b64decode(data['image/jpeg'])
                        add_output_image_slide(prs, b, title=f"Image [{idx}.{oidx}]")
                    elif 'text/plain' in data:
                        add_output_text_slide(prs, data['text/plain'], title=f"Output [{idx}.{oidx}]")
                    elif 'text/html' in data:
                        add_output_text_slide(prs, data['text/html'], title=f"HTML Output [{idx}.{oidx}]")
                # error outputs
                elif out.get('output_type') == 'error':
                    tb = "\n".join(out.get('traceback', [])) or f"{out.get('ename','')}: {out.get('evalue','')}"
                    add_output_text_slide(prs, tb, title=f"Error [{idx}.{oidx}]")

    prs.save(out_path)
    print(f"Saved presentation with outputs/images: {out_path}")

if __name__ == "__main__":
    notebook_to_pptx(NOTEBOOK_PATH, OUTPUT_PPTX)